import time, os
import io
import scipy
import numpy as np
import json
# ----------------SPEECH RECOGNITION /AUDIO/TEXT TO SPEECH DEPENDENCIES-------------#
from scipy.io.wavfile import read as wav_read
from pydub import AudioSegment
from pydub.playback import play
from mutagen.mp3 import MP3
from mutagen.wave import WAVE
import speech_recognition as sr
import ffmpeg
from gtts import gTTS
import soundfile as sf

# ----STREAMLIT RELATED----#
import streamlit as st
import requests
import re
import urllib.request
import fnmatch

# ---------DOCUMENT/WEBSITE PARSING---------#
from bs4 import BeautifulSoup
from collections import deque
from html.parser import HTMLParser
from urllib.parse import urlparse

# -------DATA FRAME/DOCX/TEXT HANDLING----------$
import pandas as pd
import pprint as pp
from pptx import Presentation
from PIL import Image
import pytesseract
from docx import Document as Docxreader
from docx.shared import Inches
from google.cloud import vision
from google.oauth2 import service_account
from google.cloud import storage
import textwrap
import glob
from spire.doc import *
from spire.doc.common import *
import fitz
from unidecode import unidecode
import pandas as pd
import numpy as np
import re
# -----------------------------------------------------------------------------------------#
# ---------------------------------OPENAI and LANGCHAIN DEPENDENCIES-----------------------#
# -----------------------------------------------------------------------------------------#
import openai
from langchain.chat_models import ChatOpenAI
from langchain.llms import OpenAI

# from langchain import HuggingFacePipeline
from langchain.llms.huggingface_pipeline import HuggingFacePipeline
from langchain.chains import RetrievalQA
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.text_splitter import CharacterTextSplitter
from langchain.document_loaders import UnstructuredWordDocumentLoader
from langchain.vectorstores import FAISS
from langchain.chains import LLMChain
from dotenv import find_dotenv, load_dotenv
from langchain.prompts.chat import (
    ChatPromptTemplate,
    SystemMessagePromptTemplate,
    HumanMessagePromptTemplate,
)
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.document_loaders import TextLoader
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate

from langchain.llms import OpenAI
from langchain import PromptTemplate, LLMChain
import random
import time
from langchain.chains import RetrievalQA
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from langchain.vectorstores import Chroma
from langchain.llms import GPT4All, LlamaCpp
from langchain.document_loaders import TextLoader, PDFMinerLoader, CSVLoader
from langchain.text_splitter import CharacterTextSplitter

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from langchain.docstore.document import Document as DocumentLang
from langchain import HuggingFacePipeline
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.embeddings import HuggingFaceInstructEmbeddings
from tqdm import tqdm
from langchain.vectorstores import FAISS
from langchain.document_loaders import (
    CSVLoader,
    EverNoteLoader,
    PDFMinerLoader,
    TextLoader,
    UnstructuredEmailLoader,
    UnstructuredEPubLoader,
    UnstructuredHTMLLoader,
    UnstructuredMarkdownLoader,
    UnstructuredODTLoader,
    UnstructuredPowerPointLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredFileLoader   
)
import os
import glob
import argparse
import openai
from time import sleep
chunk_size = 512 #512
chunk_overlap = 50
#from dotenv import load_dotenv
from langchain.chains import RetrievalQA
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from langchain.vectorstores import Chroma
from langchain.llms import GPT4All, LlamaCpp

# --------------KEYBERT MODEL ---------------#
from keybert import KeyBERT


# ----------DEFAULTS:------------#
LANGUAGE = "en"

SECRET_TOKEN = os.environ["SECRET_TOKEN"]
openai.api_key = SECRET_TOKEN
import os
os.environ["TOKENIZERS_PARALLELISM"] = "false"



# ---------------READ THE UPLOADED DOCUMENT AND GENERATE THE SPLIT---------------#
@st.cache_resource(show_spinner=True)
def readdoc_splittext(filename):
    """
    This functions takes in an input document and finds the headings, and
    splits them based on the chunks needed.
    """
    if ".docx" in filename:
        document = Docxreader(filename)
    elif (".doc" in filename) and (".docx" not in filename):
        file_path = os.path.join( os.getcwd(), "file.docx")
        document = Document()
        # Load a Word DOC file
        document.LoadFromFile(filename)
        #        # Save the DOC file to DOCX format
        document.SaveToFile(file_path, FileFormat.Docx2016)
        # Close the Document object
        document.Close()
        document = Docxreader(file_path)
    headings = []
    para_texts = []
    i = 0
    j = 0
    n = 1500  # Number of characters to be included in a single chunk of text
    t = ""
    for paragraph in document.paragraphs:
        if paragraph.style.name == "Heading 2":
            no_free_text = "".join(filter(lambda x: not x.isdigit(), paragraph.text))
            k = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ"

            getVals = list(filter(lambda x: x in k, no_free_text))
            result = "".join(getVals)
            headings.append(result.lower())
            i = 1
            # print(paragraph.text)

        elif paragraph.style.name == "normal" and i == 1:
            t += paragraph.text

            try:
                if document.paragraphs[j + 1].style.name == "Heading 2":
                    i = 0
                    # to remove numeric digits from string

                    para_texts.append(t)
                    t = ""
                    # print(i)
            except:
                # print("reached doc end")
                para_texts.append(t)
                # print(i)
        j += 1

    all_text = ""
    for text in para_texts:
        all_text += text

    a = glob.glob(filename)

    chunk_size = 1024
    chunk_overlap = 10
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=chunk_overlap
    )
    texts_isb = []
    texts_chunk = []
    documents = []
    for i in range(len(a)):
        documents.extend(UnstructuredWordDocumentLoader(a[i]).load())
        for j in range(
            len(
                text_splitter.split_documents(
                    UnstructuredWordDocumentLoader(a[i]).load()
                )
            )
        ):
            text_chunk = text_splitter.split_documents(
                UnstructuredWordDocumentLoader(a[i]).load()
            )[j]
            text_chunk.page_content = text_chunk.page_content.replace("\n", " ")
            text_chunk.page_content = text_chunk.page_content.replace("\\n", " ")
            text_chunk.page_content = text_chunk.page_content.replace("  ", " ")
            text_chunk.page_content = text_chunk.page_content.replace("  ", " ")
            texts_isb.append(text_chunk.page_content)
            texts_chunk.append(text_chunk)
    text_split = texts_isb
    return all_text, text_split, texts_chunk, headings, para_texts


def remove_newlines(serie):
    serie = serie.replace("\n", " ")
    serie = serie.replace("\\n", " ")
    serie = serie.replace("  ", " ")
    serie = serie.replace("  ", " ")
    return serie

def get_block_dict_fromDoc(doc):
    '''
    Takes in the document raw data and extracts the document text and meta data into a single dictionary.
    '''
    block_dict = {}
    page_num = 1
    for page in doc: # Iterate all pages in the document
        file_dict = page.get_text('dict') # Get the page dictionary
        block = file_dict['blocks'] # Get the block information
        block_dict[page_num] = block # Store in block dictionary
        page_num += 1 # Increase the page value by 1
    return block_dict


def get_docfeature_dataframe(block_dict):
    """
    takes in the block dict and then gives out a Pandas data frame with the most important document features.
    Once the font_size column is populated, and the index of the rows are available, one can simply order them better based on
    the font size, and document tree (sections - sub-sections - sub-sub-sections, etc)
    """
    spans = pd.DataFrame(columns=['xmin', 'ymin', 'xmax', 'ymax', 'text', 'tag'])
    rows = []
    for page_num, blocks in block_dict.items():
        for block in blocks:
            if block['type'] == 0:
                for line in block['lines']:
                    for span in line['spans']:
                        xmin, ymin, xmax, ymax = list(span['bbox'])
                        font_size = span['size']
                        text = unidecode(span['text'])
                        span_font = span['font']
                        #is_upper = False
                        is_bold = False
                        if "bold" in span_font.lower():
                            is_bold = True
                        #if font_size >= 12.0:
                        #    is_upper = True
                        if text.replace(" ","") !=  "":
                            rows.append((xmin, ymin, xmax, ymax, text, is_bold, span_font, font_size))
    span_df = pd.DataFrame(rows, columns=['xmin','ymin','xmax','ymax', 'text','is_bold','span_font', 'font_size'])
    return span_df


def get_paragraphs(headings, paragraph_sentences):

    """
    Takes the headings and the sentences extracted from the document data frame.
    Unifies the sentences under the section headings as a single paragraph.
    """
    paragraph_list = len(headings.index)*[""]
    #print(paragraph_list)
    #print(len(headings.index))

    if len(headings.index) > 1:
        #print("Length is more than 1")
        for i in range(len(headings.index)):

            paragraph = paragraph_list[i]
            #print(len(paragraph))
            if i < len(headings.index) - 1:
                #print(headings.index[i], headings.index[i+1])
                for j in range(headings.index[i]+1,headings.index[i+1]):
                    if j in paragraph_sentences.index:
                        paragraph += paragraph_sentences.loc[j]
                paragraph_list[i] = paragraph

            elif i == len(headings.index) - 1:

                for j in range(headings.index[i]+1,paragraph_sentences.index[-1]):
                    if j in paragraph_sentences.index:
                        paragraph += paragraph_sentences.loc[j]
                paragraph_list[i] = paragraph
            #print(i)

        return  paragraph_list

    elif len(headings.index) == 1:
       # print("Length is 1")
        paragraph = paragraph_list[0]
        #print(len(paragraph))

        for j in range(headings.index[0]+1, paragraph_sentences.index[-1]):
            if j in paragraph_sentences.index:
                paragraph += paragraph_sentences.loc[j]
        paragraph_list[0] = paragraph
        

        return  paragraph_list
def extract_text_from_pptx(filename):
    text = ""
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
def extract_image_addresses(filename):
    image_addresses = []

    prs = Presentation(filename)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Check if shape is an image
                # Get the image address (file path or URL)
                image_address = shape.image.filename
                # Append the image address to the list
                image_addresses.append(image_address)

    return image_addresses
#---------------READ THE .PPTX FILE AND GENERATE THE SPLIT--------------------#
@st.cache_resource(show_spinner=True)
def readdoc_splittext_pptx(filename):
    if ".pptx" in filename:
        loader = UnstructuredPowerPointLoader(filename)
        # Extract text from the PowerPoint presentation
        docs = str(loader.load())  # Assuming loader.load() returns the text content of the presentation
        
        # Extract image addresses from the PowerPoint presentation
        image_addresses = extract_image_addresses(filename)
        print(image_addresses)
        try:
            bucket_name='aiexplorers'
            os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = 'able-store-415222-3c73cfca4950.json'
            # Authenticate using service account credentials
            credentials = service_account.Credentials.from_service_account_info({
            "type": "service_account",
            "project_id": "able-store-415222",
            "private_key_id": "3c73cfca4950749ff7d8b9411e28e7db84c89c75",
            "private_key": "-----BEGIN PRIVATE KEY-----\nMIIEvgIBADANBgkqhkiG9w0BAQEFAASCBKgwggSkAgEAAoIBAQDZEH33tpJHp4XL\n3+jMloFGguZCzQuT8HwobVK7gI+x/fNfbamU2gzla+vP7GGmIP8JYn0S1cvWY0ki\nyakN7Hxp6DmkIeOpK9piojKriUFUFXu86JiaEmYRUplSeifOCexHM4KIgM4ABKc9\nogGVilJUdoLEigaMXF8/XS7O2uBRaApegUhSCZnYEp64VXVVyPoRtFAcKhmp5y4T\n+x58TjTX/7oEC6nWvfq1opJBzX++svH8f/eQ4FfpCh2WfyDyQk7GNyRxQMM9Xbnx\nwM9aWejuDN+yBufqWjoDI1wOqBNNDAuOr/Gdkpp8E5CqP/mBNwsK5LoMyWaAgrUq\nKkrE/M2fAgMBAAECggEAFXIer/m2zCOOAY2Wx5fbH+dQ7qqpiS1KZwgeE25pRBGu\nGPCRv6YnfgdCniFhqHM8H2Y+qzy6Te0zTZ+U9041IJWUFFZQ0ucFRZd24CHIdZFF\nz20/noFQ2kdJ5K8y/ZVXZdV+ARSQxcmv6FP4ShkHd1baJqbobLe3pWAEfuV7Ii+x\nlW026XwMvMZqGIn5Vw0ra4waJbexmGIHFy29pc7r8RTzgpxXEkSMi3JBPq4q6cfK\nV7ckoUw5qtvhFvhjTH+B2ipPII5+YxAo6MoxTwMrH3aZUrjTVjWzTQ3Fbb72Il77\nlKKAuzNNlNycN5T5cR3WR8rJ9lf2xfIqw9Jw44YDkQKBgQD+63n1aovkv8VLAJeJ\nbb/CRAU3ohrJpx3oRTSdwOkfVRITdJpPzxtZUun8xgFjmR50qptRNobZegbBF+FQ\nRjKu7L69UC0LEEBItn9ud+iQ64JNp+f2Wq7doTRqXHVZBANr3fSwZteQyqh2UHeF\nGQmVk7t35vIj9vXF/4sfAzhFpwKBgQDZ+/O9Y3CmyiHNX7sAJe4n//5W1VEW1nsC\nyO4+fIJPv7NPhJN8AeC//LltKKKm7C87Mc9JTkHQEtTeQoqW8/xlTwGpop068QmE\nVGWM9bfA1hLSgtwp8GoGrT92FX+12BKg5D6UxDslNDzXhrBqXGyVTqIhuPWy0zWw\nL1unWZenSQKBgQC1f2is8Dg8HMHWvcwmv+oo4Y8pZhRWgNLNXgCxVPlmIoalLX70\n1ctuFJeeLkgs5ocFn7bH4t+uFbCbaGo2YlSRsOO5HE2FtANhAfbG6z69d8Clk6eX\nAkfAapdMJxoxz6az9SrTMdXHNFMvMel27TWitrViEB9Ute+VEnW2Fe/JvwKBgAbx\nDg7+5qx5DWCD5umXS6E8drX/LwjiJaKuEWAuCNxPR3+FwkiVqrmSa7k9LQjRxqNE\n78vu0Qu2Pc8iIVWzSVtUi8ICKq2g1WPAaEd337UlXA4WGrq/LDEwPTAeeSWqTtWO\nzytfoF6L7lASuvV4IgETMviN8k0Sisgkie+nW7v5AoGBAMr4rftvEEUxOPuu+U7h\newe63BiLTU4F+OMGXhIHQIz88HdQk6w4XTj6Nrb2pivqCGPFNgTr4+WmMfe9DkAH\nxzny8Kdp9HUob6u8DrkWI2OLwKAnReVy3E3wd8MLBeZTahLHVc4vtpMH4L7cXpKv\nYlCS/pHOIfU6R5pPkSXj6ozZ\n-----END PRIVATE KEY-----\n",
            "client_email": "speekar@able-store-415222.iam.gserviceaccount.com",
            "client_id": "110531954629070375774",
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
            "auth_provider_x509_cert_url": "https://www.googleapis.com/oauth2/v1/certs",
            "client_x509_cert_url": "https://www.googleapis.com/robot/v1/metadata/x509/speekar%40able-store-415222.iam.gserviceaccount.com",
            "universe_domain": "googleapis.com"
            })

            images_text = ''
            headings_list = []
            paragraph_list = []

            # Initialize Google Cloud Vision API client
            storage_client = storage.Client()
            bucket = storage_client.get_bucket(bucket_name)
            vision_client = vision.ImageAnnotatorClient(credentials=credentials)
            feature = vision.Feature(type_=vision.Feature.Type.DOCUMENT_TEXT_DETECTION)

            # Perform OCR on images
            for i in range(len(image_addresses)):
                if os.path.exists(image_addresses[i]):  # Check if the image file exists
                    with io.open(image_addresses[i], 'rb') as image_file:
                        content = image_file.read()

                        image = vision.Image(content=content)
                        response = vision_client.text_detection(image=image)
                        texts = response.text_annotations
                        for text in texts:
                            images_text += text.description + '\n'

                        # Extract headers and paragraphs from OCR results of images
                        pat1 = re.compile(r".+\:")
                        pat2 = re.compile(r".+\.\n")
                        headings_list_img = pat1.findall(text.description)
                        paragraph_list_img = pat2.findall(text.description)

                        # Append headers and paragraphs from images to existing lists
                        headings_list.extend(headings_list_img)
                        paragraph_list.extend(paragraph_list_img)

            # Extract headers and paragraphs from the PowerPoint text
            pat1 = re.compile(r".+\:")
            pat2 = re.compile(r".+\.\n")
            headings_list_pptx = pat1.findall(docs)
            paragraph_list_pptx = pat2.findall(docs)
            headings_list.extend(headings_list_pptx)
            paragraph_list.extend(paragraph_list_pptx)
            n = 1500
            # Combine all text from PowerPoint and OCR'd images
            all_text = ''.join(paragraph_list_pptx) + images_text

            a = glob.glob(filename)

            chunk_size = 1024
            chunk_overlap = 10
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)

            texts_isb = []
            texts_raw = []
            documents = []
            for i in range(len(a)):
                documents.extend(loader.load())
                for j in range(len(text_splitter.split_documents(loader.load()))):
                    text_chunk = text_splitter.split_documents(loader.load())[j]
                    text_chunk.page_content = text_chunk.page_content.replace("\n", " ")
                    text_chunk.page_content = text_chunk.page_content.replace("\\n", " ")
                    text_chunk.page_content = text_chunk.page_content.replace("  ", " ")
                    text_chunk.page_content = text_chunk.page_content.replace("  ", " ")
                    texts_isb.append(text_chunk.page_content)
                    texts_raw.append(text_chunk)

            text_split = texts_isb

            return all_text, text_split, texts_raw, headings_list, paragraph_list
        except Exception as e:
            print("An error occurred:", e)
#---------------READ THE .TXT FILE AND GENERATE THE SPLIT--------------------#
@st.cache_resource(show_spinner=True)
def readdoc_splittext_txt(filename):
  if ".txt" in filename:
    loader = UnstructuredFileLoader(filename)
    docs = str(loader.load())
  pat1= re.compile(r".+\:")
  pat2=re.compile(r".+\.\n")
  headings_list=pat1.findall(docs)
  paragraph_list=pat2.findall(docs)
  n = 1500 #Number of characters to be included in a single chunk of text
  all_text=''
  for text in paragraph_list:
      all_text+=text
  a=glob.glob(filename)
    #print(a)
  chunk_size = 1024
  chunk_overlap=10
  text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
  texts_isb=[]
  texts_raw = []
  documents=[]
  for i in range(len(a)):
        documents.extend(UnstructuredFileLoader(a[i]).load())
        for j in range(
            len(
                text_splitter.split_documents(
                    UnstructuredFileLoader(a[i]).load()
                )
            )
        ):
            text_chunk = text_splitter.split_documents(
                UnstructuredFileLoader(a[i]).load()
            )[j]
            text_chunk.page_content = text_chunk.page_content.replace("\n", " ")
            text_chunk.page_content = text_chunk.page_content.replace("\\n", " ")
            text_chunk.page_content = text_chunk.page_content.replace("  ", " ")
            text_chunk.page_content = text_chunk.page_content.replace("  ", " ")
            texts_isb.append(text_chunk.page_content)
            texts_raw.append(text_chunk)

  text_split = texts_isb
  return all_text, text_split, texts_raw, headings_list, paragraph_list
#---------------READ THE .Pdf FILE AND GENERATE THE SPLIT--------------------#
@st.cache_resource(show_spinner=True)
def readdoc_splittext_pdf(filename):
    """
    This functi
    ons takes in an input document in pdf form and finds the headings, and
    splits them based on the chunks needed. 
    
    """
    #print("reading files")
    #print(filename)
    #print(os.path.realpath("./"), os.getcwd())
    #print(os.listdir("./"))
    #for file in os.listdir(os.getcwd()):
    #document = os.path.join("tempDir/", filename)
    doc = fitz.open(filename)
    
    
    block_dict = get_block_dict_fromDoc(doc)
    span_df = get_docfeature_dataframe(block_dict)
    doc_clean = span_df[span_df["font_size"]>=span_df["font_size"].mode()[0]]
    #doc_clean.head()
    paragraphs = doc_clean.text[doc_clean.font_size == span_df.font_size.mode()[0]]
    #print(paragraphs.values,paragraphs.index)
    headings = doc_clean.text[doc_clean.font_size > span_df.font_size.mode()[0]]
    #print(headings.values, headings.index)
    paragraph_list = get_paragraphs(headings, paragraphs)
    headings_list = headings.values.tolist()

    n = 1500 #Number of characters to be included in a single chunk of text  
    
    all_text=''
    for text in paragraph_list:
        all_text+=text

    

    a=glob.glob(filename)
    #print(a)
    chunk_size = 1024
    chunk_overlap=10
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    texts_isb=[]
    texts_raw = []
    documents=[]
    for i in range(len(a)):

        documents.extend(PDFMinerLoader(a[i]).load())
        for j in range(len(text_splitter.split_documents(PDFMinerLoader(a[i]).load()) ) ):
            text_chunk=text_splitter.split_documents(PDFMinerLoader(a[i]).load())[j]
            text_chunk.page_content = text_chunk.page_content.replace('\n', ' ')
            text_chunk.page_content = text_chunk.page_content.replace('\\n', ' ')
            text_chunk.page_content = text_chunk.page_content.replace('  ', ' ')
            text_chunk.page_content = text_chunk.page_content.replace('  ', ' ')
            # print(text_chunk)
            
            texts_isb.append(text_chunk.page_content)
            texts_raw.append(text_chunk)
    text_split=texts_isb
    return all_text, text_split,texts_raw, headings_list, paragraph_list

# ----------------CREATE CONTEXT-----------------------#
import time
from keybert import KeyBERT

def create_context(query, text_split, headings, para_texts):
    """
    Create a context for a question by finding the most similar context from the dataframe
    """

    # Create the kw_model inside the function
    kw_model = KeyBERT()

    # Get the embeddings for the entire document
    doc_keywords = kw_model.extract_keywords(" ".join(text_split))

    returns = []
    keywords_q = []

    for keyword, score in doc_keywords:
        if score > 0.3:
            keywords_q.append(keyword)

    keyword_doc = {}

    for i, sent in enumerate(text_split):
        if isinstance(sent, str) and len(sent) > 6:
            keywords = kw_model.extract_keywords(sent)
            keyword_doc_sent = []

            for keyword, score in keywords:
                if score > 0.3:
                    keyword_doc_sent.append(keyword)

            # Add the heading of the para corresponding to the sentence
            for h, pt in zip(headings, para_texts):
                pt = pt.lower()
                index = pt.find(sent.lower())
                if index != -1:
                    keyword_doc_sent.append(h)

            keyword_doc[i] = keyword_doc_sent

    returns_set = set()

    for i, keyword_list in keyword_doc.items():
        if any(keyword in keywords_q for keyword in keyword_list):
            returns_set.add(remove_newlines(text_split[i]))

    returns = list(returns_set)

    return "\n\n###\n\n".join(returns), keywords_q

def remove_newlines(s):
    return s.replace("\n", "")
"""# Usage
kw_model = KeyBERT()
context, keywords = create_context(query, text_split, headings, para_texts, kw_model)
print("Context:", context)
print("Keywords:", keywords)"""




# ------------------------SLIM KAR BASED CHATBOT----------------------------#
# ------------------------SLIM KAR BASED CHATBOT----------------------------#
#def chatbot_slim(query, context, keywords):#text_split, headings, para_texts):
#    """
#    Here, this function takes in the textual query, along with the textual context and uses KAR framework to geerate a suitable response
#    with little to almost no hallucinations. Here, openai's davinci-003 has been used to generate the response.
#    """
#
#    if input:
#        stime = time.time()
#
#        #context, keywords = create_context(query, text_split, headings, para_texts)

#       ctype = ["stuff", "map_reduce", "refine", "map_rerank"]
#      template = """
#              You are a helpful assistant who answers question based on context provided: {context}

#             If you don't have enough information to answer the question, say: "Sorry, I cannot answer that".

#              """
#        template = """
#                  You are a helpful assistant who answers question based on context provided: {context}

#                  If you don't have enough information to answer the question, say: "I cannot answer".

#                  """
#        template = """ You answer question based on context below, and if question can't be answered based on context, say \"I don't know\"\n\nContext: {context} """

#        system_message_prompt = SystemMessagePromptTemplate.from_template(template)

#        # Human question prompt

#        human_template = "Answer following question: {question}"

#        template = """ Answer question {question} based on context below, and if question can't be answered based on context,
#       say \"I don't know\"\n\nContext: {context}

#        Answer:
#        """

#        template = """ Use following pieces of context to answer the question. Provide answer in full detail using provided context.
#        If you don't know the answer, say I don't know
#        {context}
#        Question : {question}
#        Answer:"""

#        human_message_prompt = HumanMessagePromptTemplate.from_template(human_template)

#        chat_prompt = ChatPromptTemplate.from_messages(
#            [system_message_prompt, human_message_prompt]
#        )

#        chunk_size = 1500
#        PROMPT = PromptTemplate(
#            input_variables=["context", "question"], template=template
#        )
#
#        chain_type_kwargs = {"prompt": PROMPT}

#        question = query
#        context = context
#        openai.api_key = SECRET_TOKEN
#        model = "gpt-3.5-turbo-instruct"
#        chat = openai.Completion.create(
#            #prompt=f"You answer question based on context below, and if the question can't be answered based on the context, 
#            #say \"I don't know\"\n\nContext: {context}\n\n---\n\nQuestion: {question}\nAnswer:",
#            prompt=f"""You are a question answering assistant with no previous information, 
#            you answer question based on following context and if question cannot be answered based on context, 
#            say \"I don't know\"\n\nContext: {context}\n\n---\n\nQuestion: {question}\nAnswer:""",
#            temperature=0,
#            max_tokens=2000,
#            top_p=1,
#            frequency_penalty=0,
#            presence_penalty=0,
#            stop=None,
#            model=model,
#        )

#        reply = chat["choices"][0][
#            "text"
#        ].strip()  # ['choices'][0]['message']['content']

#        return reply, context, keywords
from keybert import KeyBERT
import openai

def chatbot_slim(question, context, keywords):
    """
    This function takes in the textual query, along with the textual context and uses KAR framework to generate a suitable response
    with little to almost no hallucinations. Here, openai's davinci-003 has been used to generate the response.
    """

    # Create the kw_model inside the function
    kw_model = KeyBERT()

    # Get the embeddings for the entire document
    doc_keywords = kw_model.extract_keywords(" ".join(context))

    returns = []
    keywords_q = []

    for keyword, score in doc_keywords:
        if score > 0.3:
            keywords_q.append(keyword)

    keyword_doc = {}

    for i, sent in enumerate(context):
        if isinstance(sent, str) and len(sent) > 6:
            keywords = kw_model.extract_keywords(sent)
            keyword_doc_sent = []

            for keyword, score in keywords:
                if score > 0.3:
                    keyword_doc_sent.append(keyword)

            keyword_doc[i] = keyword_doc_sent

    returns_set = set()

    for i, keyword_list in keyword_doc.items():
        if any(keyword in keywords_q for keyword in keyword_list):
            returns_set.add(sent)

    returns = list(returns_set)

    # Create conversation prompt
    conversation_prompt = ""
    # Add the user's question to the conversation
    conversation_prompt += f"User: {question}\nAssistant:"
    # Adding recorded conversation to the prompt
    for i, user_input in enumerate(returns):
        conversation_prompt += f"User: {user_input}\nAssistant:"
        if i < len(returns) - 1:
            conversation_prompt += "\n"

    # Add the user's question to the conversation
    #conversation_prompt += f"User: {question}\nAssistant:"

    openai.api_key = "YOUR_SECRET_TOKEN"  # Replace with your actual secret token
    model = "gpt-3.5-turbo-0125"

    # Generate response using OpenAI's Chat API
    chat = openai.ChatCompletion.create(
        model=model,
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": conversation_prompt},
        ],
        temperature=0,
        max_tokens=2000,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0,
        stop=None,
    )

    response = chat["choices"][0]["message"]["content"].strip()

    # Return the response and context
    return response, context, keywords



@st.cache_resource(show_spinner=True)
def create_db(texts_raw,_uploaded_file_name):


    hf= OpenAIEmbeddings(model="text-embedding-ada-002", openai_api_key=openai.api_key)
    db = FAISS.from_documents(texts_raw, hf)
    db.save_local("faiss_index_anupam" + _uploaded_file_name)
    db=FAISS.load_local("faiss_index_anupam" + _uploaded_file_name, hf)
    return hf, db
    
def chatbot(question, db):

    openai.api_key = SECRET_TOKEN
    
    ctype=['stuff', 'map_reduce', 'refine', 'map_rerank']
    st.write(question)

    retriever = db.as_retriever(search_type='similarity', search_kwargs={"k": 4} )#do not increase k beyond 3, else
    docs_and_scores = db.similarity_search_with_score(question)
        
    
    llm = OpenAI(model='gpt-3.5-turbo-instruct',temperature=0, openai_api_key=openai.api_key)
    qa = RetrievalQA.from_chain_type(llm=llm, chain_type=ctype[0], retriever=retriever, return_source_documents=True)

    
    query = question
    
    res = qa(query)
    response = openai.ChatCompletion.create(
    model="gpt-3.5-turbo-0125",
    
    messages=[
        {"role": "system", 
         "content": "You are a question-answering assistant, with no prior information. You answer questions based on content provided, if you dont know the answer say I do not know"},  
        {"role": "user", "content": f"{res}"}
    ])
    answer= response["choices"][0]["message"]["content"]
    return answer


# ----------------TEXT TO SPEECH FUNCTION FOR ANSWER READOUT---------#
def texttospeech_raw(text, language, savename="answer", slow=False):
    """
    This function here, calls the google text to speech engine to read out the answer generated by the KAR framework.
    """
    myobj = gTTS(text=text, lang=language, slow=False)

    # Saving the converted audio in a mp3 file
    myobj.save(savename + ".mp3")
    sound = AudioSegment.from_mp3(savename + ".mp3")
    sound.export(savename + ".wav", format="wav")


# ---------------SPEECH RECOGNITION---------#
def speechtotext(query_audio):
    """
    This function takes in a ".wav" audio file as input and converts into text.
    """
    r = sr.Recognizer()

    audio_ex = sr.AudioFile(query_audio)
    # type(audio_ex)

    # Create audio data
    with audio_ex as source:
        audiodata = r.record(audio_ex)
    # type(audiodata)
    # Extract text
    text = r.recognize_google(audio_data=audiodata, language="en-US")
    #print("stotext ::", text)
    return text
